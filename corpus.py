import os
import glob
import pandas as pd
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract
import json
from odf.opendocument import load
from odf.text import P
from pdf2image import convert_from_path

# Folder path
folder_path = ""
os.makedirs(folder_path, exist_ok=True)

# Collect all files
files = glob.glob(os.path.join(folder_path, '**'), recursive=True)

# Text corpus
corpus = ""

def add_delimiter(filename):
    return f"\n\n---\nFILE: {filename}\n---\n\n"

# --- File processors ---

def process_pdf(file):
    text_chunks = []
    reader = PdfReader(file)

    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if text.strip():
            text_chunks.append(text)
        else:
            # OCR fallback
            images = convert_from_path(file, first_page=i+1, last_page=i+1)
            for img in images:
                ocr_text = pytesseract.image_to_string(img, lang="eng+heb")
                if ocr_text.strip():
                    text_chunks.append(ocr_text)
    return "\n".join(text_chunks)

def process_csv(file):
    return pd.read_csv(file).to_string()

def process_parquet(file):
    return pd.read_parquet(file).to_string()

def process_docx(file):
    doc = Document(file)
    return "\n".join(para.text for para in doc.paragraphs)

def process_odt(file):
    doc = load(file)
    return "\n".join(p.firstChild.data for p in doc.getElementsByType(P) if p.firstChild)

def process_pptx(file):
    prs = Presentation(file)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)

def process_image(file):
    return pytesseract.image_to_string(Image.open(file), lang="eng+heb")

def process_code(file):
    with open(file, 'r', encoding="utf-8", errors="ignore") as f:
        return f.read()

def process_ipynb(file):
    with open(file, 'r', encoding="utf-8") as f:
        notebook = json.load(f)
    text_blocks = []
    for cell in notebook.get("cells", []):
        if cell["cell_type"] in ("markdown", "code"):
            text_blocks.append("".join(cell["source"]))
    return "\n".join(text_blocks)

def process_video(file):
    return f"Video file (not parsed): {file}"

# --- Dispatcher ---
processors = {
    (".pdf",): process_pdf,
    (".csv",): process_csv,
    (".parquet",): process_parquet,
    (".docx",): process_docx,
    (".odt",): process_odt,
    (".pptx",): process_pptx,
    (".png", ".jpg", ".jpeg"): process_image,
    (".py", ".js", ".html", ".css"): process_code,
    (".ipynb",): process_ipynb,
    (".mp4", ".avi", ".mkv"): process_video,
}

# --- Process all files ---
for file in files:
    if not os.path.isfile(file):
        continue
    ext = os.path.splitext(file)[1].lower()
    for exts, handler in processors.items():
        if ext in exts:
            print(f"Processing file: {file}")
            try:
                content = handler(file)
                if content.strip():
                    corpus += add_delimiter(file)
                    corpus += content.strip() + "\n"
            except Exception as e:
                print(f"Error processing {file}: {e}")
            break

# --- Save result ---
corpus_path = os.path.join(folder_path, "corpus.txt")
with open(corpus_path, "w", encoding="utf-8") as f:
    f.write(corpus)

print(f"\n Consolidated corpus saved at {corpus_path}")
print(f"Total characters: {len(corpus)}")
